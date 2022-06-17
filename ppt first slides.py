from datetime import datetime
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import statistics
from datetime import date, timedelta, datetime
import io
from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation



# Import data
df=pd.read_csv('concat_datas_elea.csv', index_col=0, sep=';')
dictio = pd.read_csv('Dict_final.csv', delimiter=';')

# Create ppt
ppt = Presentation()

# First slide
group=df.groupby(['ehr'])
title_slide_layout = ppt.slide_layouts[0]
slide = ppt.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "CLARIFY"
subtitle = slide.placeholders[1]
subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())

# Second slide
# Initial analysis
info=[]
for i in df.columns:
    info.append(f'{i} : {df[i].nunique()} values')

second_slide = ppt.slide_layouts[5]
second = ppt.slides.add_slide(second_slide)
second.shapes.title.text = "Data has been extracted"
# Add a text box
textbox = second.shapes.add_textbox(Inches(3), Inches(1.5),Inches(3), Inches(1))
textframe = textbox.text_frame
paragraph = textframe.add_paragraph()
paragraph.text = "Data from " + str(len(df)) +" women has been processed"
# Add another text box
textbox2 = second.shapes.add_textbox(Inches(1), Inches(2),Inches(1), Inches(10))
textframe2 = textbox2.text_frame
paragraph2 = textframe2.add_paragraph()
paragraph2.text =  "Unique values:  "
#info[1:(round(len(info)/2))]


# Third slide
third_slide = ppt.slide_layouts[2]
third = ppt.slides.add_slide(third_slide)
third.shapes.title.text = "Statistics"
# table to add
table = {}
for i in df.columns:
    table[i] = df[i].nunique()

df_table = pd.DataFrame.from_dict(table, orient='index')

x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
shape = third.shapes.add_table(round(len(df_table)/2), 2, x, y, cx, cy)

#table_placeholder = third.shapes[1]
#shape = table_placeholder.insert_table(rows=3, cols=4)


fourth = ppt.slides.add_slide(ppt.slide_layouts[2])
shape = fourth.shapes.add_table(len(df_table)+1, 2, x, y, cx, cy)

table = shape.table

cell = table.cell(0, 1)
cell.text
cell.text = "Number of patients with data"

row=1

for i in df_table.index:
    cell = table.cell(row, 0)
    cell.text
    cell.text = str(i)
    row+=1

row2=1
for i in df_table.iloc[:,0]:
    cell = table.cell(row2, 1)
    cell.text
    cell.text = str(i)
    row2 += 1


ppt.save("test.pptx")