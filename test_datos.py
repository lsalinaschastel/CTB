
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import statistics
from statistics import mode
from datetime import date, timedelta, datetime
import io
from pptx import Presentation
from pptx.util import Inches



# Test data
#df = pd.read_csv('COVID19_data.csv', index_col=0)
#dictio = pd.read_csv('Dict.csv', delimiter=';')

# Real data
df=pd.read_csv('concat_datas_elea.csv', index_col=0, sep=';')
#dictio=pd.read_excel('Dict_final.xlsx',index_col=None)
#df = pd.read_csv('concat_datas.csv', index_col=0, delimiter=';')
dictio = pd.read_csv('Dict_final_elea.csv', delimiter=';')

#df=pd.read_csv('concat_data_final.csv', index_col=0, sep=';')
#dictio=pd.read_excel('Dict_final.xlsx',index_col=None)

#Fill in missing values
#df=df.fillna('Missing values')
# Remove date columns (if it does not say in the type it is a date)
# for i in df.columns:
#     if df[i].dtype=='datetime64[ns]':
#         df=df.drop([i],axis=1)
#         dictio=dictio.drop([i])


# Create ppt
ppt = Presentation()

# Remove columns that have no values



# Separate cat from numerical variables
cat_vars = []
num_vars = []
date_vars=[]

for i, row in dictio.iterrows():
    if dictio['Type'][i] == 'num':
        num_vars.append(dictio['Variable'][i])
    elif dictio['Type'][i] == 'cat':
        cat_vars.append(dictio['Variable'][i])
        # cat_vars.append(dictio.index[i])
    elif dictio['Type'][i] == 'date':
        date_vars.append(dictio['Variable'][i])


df_cat = df[cat_vars]
df_num = df[num_vars]
df_dates = df[date_vars]

# Analysis of DATES variables

# Convert all columns to date type

for col in df_dates.columns:
    if df_dates[col].dtypes == 'O':
        df_dates[col] = pd.to_datetime(df_dates[col])

    if df_dates[col].dtypes == 'float64' or 'int64':
        df_dates[col] = pd.to_datetime(df_dates[col], format='%Y')

# Finding columns "automatically"
dx_date=df_dates.filter(regex='diagnosis|dx').columns
first_treat=df_dates.filter(regex='first_treat').columns
rec_year=df_dates.filter(regex='recurrence|rec').columns
death_year=df_dates.filter(regex='death_year|death_date').columns

# Operations with dx_date
days_dx_treat=df_dates[first_treat].iloc[:,0]-df_dates[dx_date].iloc[:,0]
days_dx_surg1=df_dates['date_1']-df_dates[dx_date].iloc[:,0]
#days_dx_rec=df_dates[rec_year].iloc[:,0]-df_dates[dx_date].iloc[:,0]
days_dx_death=df_dates[death_year].iloc[:,0]-df_dates[dx_date].iloc[:,0]


# Plots
sns.distplot(days_dx_treat.dt.days) #scatter plot
plt.title('Days until treatment')
plt.ylabel("frequency")
plt.xlabel("days")
#plt.xlim([0,200])
plt.savefig("Dates1", dpi=300)

# Pict to ppt
img_path1 = "Dates1.png"
graph_slide_layout = ppt.slide_layouts[8]
slide = ppt.slides.add_slide(graph_slide_layout)
placeholder = slide.placeholders[1]
pic = placeholder.insert_picture(img_path1)


sns.distplot(days_dx_surg1.dt.days) #scatter plot
plt.title('Days until first surgery')
plt.ylabel("frequency")
plt.xlabel("days")
#plt.xlim([0,200])
plt.savefig("Dates2", dpi=300)

# Pict to ppt
img_path2 = "Dates2.png"
graph_slide_layout = ppt.slide_layouts[8]
slide = ppt.slides.add_slide(graph_slide_layout)
placeholder = slide.placeholders[1]
pic = placeholder.insert_picture(img_path2)


sns.distplot(days_dx_treat.dt.days) #scatter plot por qué sale asi? Aún haciendo dropna da la freq tan baja.
plt.title('Days until death')
plt.ylabel("frequency")
plt.xlabel("days")
#plt.xlim([0,200])
plt.savefig("Dates3", dpi=300)

# Pict to ppt
img_path3 = "Dates3.png"
graph_slide_layout = ppt.slide_layouts[8]
slide = ppt.slides.add_slide(graph_slide_layout)
placeholder = slide.placeholders[1]
pic = placeholder.insert_picture(img_path3)




# Analysis of CATEGORICAL variables
# si están todas las columnas vacías quitar datos.
df_cat=df_cat.fillna('Missing values')
for i in df_cat.columns:

    if len(df_cat[i].value_counts()) <= 2:
        # Pie chart
        fig = plt.figure(figsize=(20, 5))
        df_cat.groupby(i).size().plot(kind='pie', textprops={'fontsize': 15},
                                      colors=['gold', 'blue'], autopct=lambda x: str(round(x, 2)) + '%',
                                      pctdistance=0.5)


        plt.ylabel("")
        plt.savefig("Piechart_variable_" + str(i), dpi=300)

        # Pict to ppt
        img_path = "Piechart_variable_" + str(i) + '.png'

        graph_slide_layout = ppt.slide_layouts[8]
        slide = ppt.slides.add_slide(graph_slide_layout)

        title = slide.shapes.title
        title.text = str(i)
        placeholder = slide.placeholders[1]
        pic = placeholder.insert_picture(img_path)
        subtitle = slide.placeholders[2]
        subtitle.text = str(len(df_cat[i].value_counts())) + " categories" + "\nSample size: " + str(len(df_cat)) +  " (" + str(round(((df_cat[i].isnull().sum() / df_cat.shape[0]) * 100),2)) + " %)" + "\nMode: " + str(statistics.mode(df_cat[i]))

    else:
        # Bar chart
        fig = plt.figure()
        df_cat.groupby(i).size().plot(kind='bar', rot=0)
        plt.savefig("Barchart_variable_" + str(i), dpi=300)

        # Pict to ppt
        img_path = "Barchart_variable_" + str(i) + '.png'
        graph_slide_layout = ppt.slide_layouts[8]
        slide = ppt.slides.add_slide(graph_slide_layout)
        title = slide.shapes.title
        title.text = str(i)
        placeholder = slide.placeholders[1]
        pic = placeholder.insert_picture(img_path)
        subtitle = slide.placeholders[2]
        subtitle.text = str(len(df_cat[i].value_counts())) + " categories" + "\nSample size: " + str(
            len(df_cat))  + " (" + str(
            round(((df_cat[i].isnull().sum() / df_cat.shape[0]) * 100), 2)) + " %)" + "\nMode: " + str(statistics.mode(df_cat[i]))






# Analysis of QUANTITATIVE variables

#Delete variables with no values
# Find the columns where each value is null
cols_empty = [col for col in df_num.columns if df_num[col].isnull().all()]
# Drop these columns from the dataframe
df_num.drop(cols_empty,
        axis=1,
        inplace=True)


for i in df_num.columns:
    print(i)
    # Box plot
    fig = plt.figure()
    plt.subplot(1, 2, 1)
    df_num[i].dropna(axis=0).plot.box(fontsize=8)
    plt.title("Box plot")

    # Finding outliers
    outlier_free_list = []
    outliers = []
    perc_outliers = []
    Q3 = np.nanquantile(df_num[i], 0.75)
    Q1 = np.nanquantile(df_num[i], 0.25)

    IQR = Q3 - Q1

    lower_range = Q1 - 1.5 * IQR
    upper_range = Q3 + 1.5 * IQR

    # Filter the data that are free of outliers
    for x in df_num[i]:
        if (x > lower_range) & (x < upper_range):
            outlier_free_list.append(x)

        else:
            outliers.append(x)
    perc_outliers.append((len(outliers) / len(df_num)) * 100)


    # Histogram
    plt.subplot(1, 2, 2)
    sns.distplot(df_num[i]) #scatter plot
    plt.title('Distribution')
    plt.ylabel("")
    plt.savefig("Num_variable_" + str(i), dpi=300)

    # Picts to ppt
    img_path = "Num_variable_" + str(i) + '.png'

    graph_slide_layout = ppt.slide_layouts[8]
    slide = ppt.slides.add_slide(graph_slide_layout)

    title = slide.shapes.title
    title.text = str(i)
    placeholder = slide.placeholders[1]
    pic = placeholder.insert_picture(img_path)
    subtitle = slide.placeholders[2]
    subtitle.text = "Variable " + str(i) + " has: " + str(len(outliers)) + " outliers" + "\nMean: " + str(round(np.nanmean(df_num[i]),2)) + "\nMedian: " + str(np.nanmedian(df_num[i])) + "\nMissing values: " + str(df_num[i].isnull().sum())




ppt.save("Test.pptx")