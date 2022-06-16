from datetime import datetime
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import statistics
from datetime import date, timedelta, datetime
import io
from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation



# Import data
df=pd.read_csv('concat_datas_elea.csv', index_col=0, sep=';')
dictio = pd.read_csv('Dict_final.csv', delimiter=';')

# Create ppt
ppt = Presentation()

# First slide
title_slide_layout = ppt.slide_layouts[0]
slide = ppt.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "CLARIFY"
subtitle = slide.placeholders[1]
subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())

# Second slide
# Initial analysis
info=[]
for i in df.columns:
    info.append(f'{i} : {df[i].nunique()} values')

second_slide = ppt.slide_layouts[5]
second = ppt.slides.add_slide(second_slide)
second.shapes.title.text = "Data has been extracted"  + "\n Unique values:  " +"\n " +str(info)
# add a text box
textbox = second.shapes.add_textbox(Inches(3), Inches(1.5),Inches(3), Inches(1))
textframe = textbox.text_frame
paragraph = textframe.add_paragraph()
paragraph.text = "Data from " + str(len(df)) +" women has been processed"




# Separate cat from numerical variables
cat_vars = []
num_vars = []
date_vars = []
for i, row in dictio.iterrows():
    if dictio['Type'][i] == 'num':
        num_vars.append(dictio['Variable'][i])
    elif dictio['Type'][i] == 'cat':
        cat_vars.append(dictio['Variable'][i])
    #elif dictio['Type'][i] == 'date':
        #date_vars.append(dictio['Variable'][i])


df_cat = df[cat_vars]
df_num = df[num_vars]
df_date = df[date_vars]


df = df.fillna('Missing values')

# Analysis of CATEGORICAL variables
#df_cat = df_cat.fillna('Missing values')
for i in df_cat.columns:
    if df_cat[i].dtypes == 'O':
        df_cat[i] = df_cat[i].str.lower()

    if len(df_cat[i].value_counts()) <= 2:
        # Pie chart
        fig = plt.figure(figsize=(20, 5))
        df_cat.groupby(i).size().plot(kind='pie', textprops={'fontsize': 10},
                                      colors=['gold', 'blue'], autopct=lambda x: str(round(x, 2)) + '%',
                                      pctdistance=0.5)
        plt.ylabel("")
        plt.savefig("Piechart_variable_" + str(i), dpi=300)

        # Pict to ppt
        img_path = "Piechart_variable_" + str(i) + '.png'
        graph_slide_layout = ppt.slide_layouts[8]
        slide = ppt.slides.add_slide(graph_slide_layout)
        title = slide.shapes.title
        title.text = str(i)
        placeholder = slide.placeholders[1]
        pic = placeholder.insert_picture(img_path)
        subtitle = slide.placeholders[2]
        subtitle.text = str(len(df_cat[i].value_counts())) + " categories" + "\nSample size: " + str(len(df_cat)) +  "\nMode: " + str(statistics.mode(df_cat[i])) + "\nMissing values: " + str(df_cat[i].isnull().sum()) +" (" + str(round(((df_cat[i].isnull().sum() / df_cat.shape[0]) * 100),2)) + " %)"

    else:
        # Bar chart
        fig = plt.figure()
        df_cat.groupby(i).size().plot(kind='barh')
        plt.xlabel("Frequency")
        plt.ylabel("")
        plt.tight_layout()
        plt.savefig("Barchart_variable_" + str(i), dpi=300)


        # Pict to ppt
        img_path = "Barchart_variable_" + str(i) + '.png'
        graph_slide_layout = ppt.slide_layouts[8]
        slide = ppt.slides.add_slide(graph_slide_layout)
        title = slide.shapes.title
        title.text = str(i)
        placeholder = slide.placeholders[1]
        pic = placeholder.insert_picture(img_path)
        subtitle = slide.placeholders[2]
        subtitle.text = str(len(df_cat[i].value_counts())) + " categories" + "\nSample size: " + str(
            len(df_cat))  +  "\nMode: " + str(statistics.mode(df_cat[i])) + "\nMissing values: " + str(df_cat[i].isnull().sum()) +" (" + str(round(((df_cat[i].isnull().sum() / df_cat.shape[0]) * 100),2)) + " %)"





# Analysis of QUANTITATIVE variables

# Delete variables with no values. Find the columns where each value is null
cols_empty = [col for col in df_num.columns if df_num[col].isnull().all()]
# Drop these columns from the dataframe
df_num.drop(cols_empty,
        axis=1,
        inplace=True)

for i in df_num.columns:
    print(i)
    # Box plot
    fig = plt.figure()
    plt.subplot(1, 2, 1)
    df_num[i].dropna(axis=0).plot.box(fontsize=8)
    plt.title("Box plot")

    # Finding outliers
    outlier_free_list = []
    outliers = []
    perc_outliers = []
    Q3 = np.nanquantile(df_num[i], 0.75)
    Q1 = np.nanquantile(df_num[i], 0.25)

    IQR = Q3 - Q1

    lower_range = Q1 - 1.5 * IQR
    upper_range = Q3 + 1.5 * IQR

    # Filter the data that are free of outliers
    for x in df_num[i]:
        if (x > lower_range) & (x < upper_range):
            outlier_free_list.append(x)

        else:
            outliers.append(x)
    perc_outliers.append((len(outliers) / len(df_num)) * 100)


    # Histogram
    plt.subplot(1, 2, 2)
    sns.distplot(df_num[i]) #scatter plot
    plt.title('Distribution')
    plt.ylabel("")
    plt.savefig("Num_variable_" + str(i), dpi=300)

    # Picts to ppt
    img_path = "Num_variable_" + str(i) + '.png'
    graph_slide_layout = ppt.slide_layouts[8]
    slide = ppt.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = str(i)
    placeholder = slide.placeholders[1]
    pic = placeholder.insert_picture(img_path)
    subtitle = slide.placeholders[2]
    subtitle.text = "Variable " + str(i) + " has: " + str(len(outliers)) + " outliers" + "\nMean: " + str(round(np.nanmean(df_num[i]),2)) + "\nMedian: " + str(np.nanmedian(df_num[i])) + "\nMissing values: " + str(df_num[i].isnull().sum())




# Plot the dates


ppt.save("Datas_final.pptx")